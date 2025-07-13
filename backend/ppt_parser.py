from pptx import Presentation
import re
from typing import List, Dict, Any
import os

class PPTParser:
    def __init__(self):
        # Regex patterns for different number formats
        self.number_patterns = [
            r'₹\s*[\d,]+\.?\d*\s*[KkMmBbCrLacslacs]*',  # Indian currency
            r'\$\s*[\d,]+\.?\d*\s*[KkMmBb]*',  # USD
            r'[\d,]+\.?\d*\s*%',  # Percentages
            r'[\d,]+\.?\d*\s*[KkMmBbCrLacslacs]+',  # Numbers with units
            r'[\d,]+\.?\d*',  # Plain numbers
        ]
        
    def parse_presentation(self, ppt_path: str) -> List[Dict[str, Any]]:
        """Parse PowerPoint presentation and extract slides with numbers"""
        try:
            prs = Presentation(ppt_path)
            slides_data = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "title": self._extract_slide_title(slide),
                    "numbers": [],
                    "text_content": []
                }
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        slide_data["text_content"].append(text_content)
                        
                        # Find numbers in text
                        numbers = self._extract_numbers_from_text(text_content)
                        for number_info in numbers:
                            number_info["slide_number"] = slide_idx + 1
                            number_info["context"] = text_content[:100] + "..." if len(text_content) > 100 else text_content
                            slide_data["numbers"].append(number_info)
                
                # Extract from tables if present
                table_numbers = self._extract_from_tables(slide, slide_idx + 1)
                slide_data["numbers"].extend(table_numbers)
                
                slides_data.append(slide_data)
            
            print(f"Parsed {len(slides_data)} slides from presentation")
            return slides_data
            
        except Exception as e:
            print(f"Error parsing presentation: {str(e)}")
            raise Exception(f"Failed to parse presentation: {str(e)}")
    
    def _extract_slide_title(self, slide) -> str:
        """Extract title from slide"""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Assume first text shape is title
                return shape.text.strip()
        return f"Slide {slide.slide_id}"
    
    def _extract_numbers_from_text(self, text: str) -> List[Dict[str, Any]]:
        """Extract numbers and their context from text"""
        numbers = []
        
        for pattern in self.number_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                number_text = match.group().strip()
                
                # Get surrounding context
                start = max(0, match.start() - 50)
                end = min(len(text), match.end() + 50)
                context = text[start:end].strip()
                
                # Parse the number
                parsed_number = self._parse_number(number_text)
                
                if parsed_number is not None:
                    numbers.append({
                        "raw_text": number_text,
                        "parsed_value": parsed_number,
                        "context": context,
                        "position": match.start(),
                        "type": self._classify_number_type(number_text)
                    })
        
        return numbers
    
    def _extract_from_tables(self, slide, slide_number: int) -> List[Dict[str, Any]]:
        """Extract numbers from tables in slide"""
        numbers = []
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        if cell.text.strip():
                            cell_numbers = self._extract_numbers_from_text(cell.text)
                            for number_info in cell_numbers:
                                number_info["slide_number"] = slide_number
                                number_info["table_position"] = f"Row {row_idx + 1}, Col {col_idx + 1}"
                                numbers.append(number_info)
        
        return numbers
    
    def _parse_number(self, number_text: str) -> float:
        """Parse number text to float value"""
        try:
            # Remove currency symbols and spaces
            clean_text = re.sub(r'[₹$\s]', '', number_text)
            
            # Handle units
            multiplier = 1
            if re.search(r'[Cc]r', clean_text, re.IGNORECASE):
                multiplier = 10000000  # 1 Crore = 10 Million
                clean_text = re.sub(r'[Cc]r.*', '', clean_text, flags=re.IGNORECASE)
            elif re.search(r'[Ll]ac', clean_text, re.IGNORECASE):
                multiplier = 100000  # 1 Lakh = 100,000
                clean_text = re.sub(r'[Ll]ac.*', '', clean_text, flags=re.IGNORECASE)
            elif re.search(r'[Kk]', clean_text):
                multiplier = 1000
                clean_text = re.sub(r'[Kk].*', '', clean_text)
            elif re.search(r'[Mm]', clean_text):
                multiplier = 1000000
                clean_text = re.sub(r'[Mm].*', '', clean_text)
            elif re.search(r'[Bb]', clean_text):
                multiplier = 1000000000
                clean_text = re.sub(r'[Bb].*', '', clean_text)
            
            # Handle percentage
            is_percentage = '%' in clean_text
            clean_text = clean_text.replace('%', '')
            
            # Remove commas and parse
            clean_text = clean_text.replace(',', '')
            number = float(clean_text)
            
            # Apply multiplier
            number *= multiplier
            
            # Convert percentage to decimal if needed for calculations
            if is_percentage:
                return number  # Keep as percentage for display
            
            return number
            
        except (ValueError, AttributeError):
            return None
    
    def _classify_number_type(self, number_text: str) -> str:
        """Classify the type of number"""
        if '₹' in number_text or '$' in number_text:
            return "currency"
        elif '%' in number_text:
            return "percentage"
        elif re.search(r'[KkMmBbCrLac]', number_text, re.IGNORECASE):
            return "metric"
        else:
            return "number"